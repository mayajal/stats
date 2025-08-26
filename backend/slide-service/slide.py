import base64
import io
import os
import pandas as pd
from flask import Flask, jsonify, request, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import logging
import json
import re # Import re module

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": ["https://vita.chloropy.com","https://vita-zur5dejluq-uc.a.run.app", "http://localhost:9002", "https://vita--statviz-j3txi.us-central1.hosted.app"]}})

# Health Check endpoint for Cloud Run
@app.route("/")
def health():
    return "OK", 200

def add_table_to_slide(slide, df, left, top, width, height):
    """Adds a pandas DataFrame as a table to a slide."""
    rows, cols = df.shape[0] + 1, df.shape[1]
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table

    # Set column headers
    for col_index, col_name in enumerate(df.columns):
        cell = table.cell(0, col_index)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add data rows
    for row_index, (idx, row_data) in enumerate(df.iterrows()):
        for col_index, cell_data in enumerate(row_data):
            cell_text = str(cell_data) if not pd.isna(cell_data) else ""
            table.cell(row_index + 1, col_index).text = cell_text

    # Adjust font size
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10)

# Helper function to clean factor names
def clean_factor_name(name):
    q_matches = re.findall(r"Q\([\'\"]([^\'\"]*)[\'\"]\)", name)
    if q_matches:
        return ":".join(sorted(list(set(q_matches))))
    return name

@app.route('/slide/generate', methods=['POST'])
def generate_slide():
    """
    Generates a PowerPoint with 2 slides per run:
    Slide 1: Left panel (Normality + ANOVA), Right panel (Mean separation table)
    Slide 2: Four equal panels, each with a chart (up to 4 charts)
    """
    if 'analysis_title' not in request.form or 'json_data' not in request.form:
        return jsonify({"error": "Missing 'analysis_title' or 'json_data' in the request form."}), 400

    analysis_title = request.form.get('analysis_title')
    json_data_str = request.form.get('json_data')

    try:
        data = json.loads(json_data_str)
    except Exception:
        logging.error("Error: Invalid JSON data provided.")
        return jsonify({"error": "Invalid input."}), 400

    try:
        logging.info("Starting slide generation.")
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_slide_layout = prs.slide_layouts[6]

        # --- Slide 1 ---
        slide1 = prs.slides.add_slide(blank_slide_layout)
        # Left panel (Normality + ANOVA)
        left_panel_left = Inches(0.5)
        left_panel_top = Inches(0.5)
        left_panel_width = Inches(6.5)
        # Title
        title_shape = slide1.shapes.add_textbox(left_panel_left, left_panel_top, left_panel_width, Inches(0.5))
        title_shape.text_frame.text = analysis_title
        title_shape.text_frame.paragraphs[0].font.size = Pt(24)
        title_shape.text_frame.paragraphs[0].font.bold = True
        # Normality Test
        shapiro_stat = data.get('shapiro', {}).get('stat')
        shapiro_p = data.get('shapiro', {}).get('p')
        stat_text = f"{shapiro_stat:.4f}" if isinstance(shapiro_stat, (int, float)) else "N/A"
        p_text = f"{shapiro_p:.4f}" if isinstance(shapiro_p, (int, float)) else "N/A"
        normality_text = f"Normality Test (Shapiro-Wilk): Statistic={stat_text}, p-value={p_text}"
        normality_shape = slide1.shapes.add_textbox(left_panel_left, Inches(1.0), left_panel_width, Inches(1.0))
        text_frame = normality_shape.text_frame
        p1 = text_frame.paragraphs[0]
        p1.text = normality_text
        p1.font.size = Pt(12)
        if isinstance(shapiro_p, (int, float)):
            if shapiro_p > 0.05:
                interpretation_text = "Residuals appear to be normally distributed (p >= 0.05)."
            else:
                interpretation_text = "Residuals appear NOT to be normally distributed (p <= 0.05)."
            p2 = text_frame.add_paragraph()
            p2.text = interpretation_text
            p2.font.size = Pt(10)
        # ANOVA Table
        anova_table_json = data.get('anova_table')
        if anova_table_json:
            anova_title_shape = slide1.shapes.add_textbox(left_panel_left, Inches(1.75), left_panel_width, Inches(0.2))
            anova_title_shape.text_frame.text = "ANOVA Table"
            p = anova_title_shape.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.LEFT
            try:
                anova_df = pd.read_json(io.StringIO(anova_table_json), orient='columns')
            except (ValueError, TypeError):
                anova_df = pd.read_json(anova_table_json)
            anova_df.reset_index(inplace=True)
            anova_df.rename(columns={'index': 'Source of Error'}, inplace=True)
            anova_df['Source of Error'] = anova_df['Source of Error'].apply(clean_factor_name)
            if 'PR(>F)' in anova_df.columns:
                # Format PR(>F) to 4 decimals (as string, keep null as is)
                anova_df['PR(>F)'] = anova_df['PR(>F)'].apply(lambda x: f"{x:.4f}" if pd.notnull(x) else x)
                def interpret_pval(pval):
                    if pd.isnull(pval):
                        return None
                    elif float(pval) <= 0.05:
                        return 'Significant'
                    else:
                        return 'Non-significant'
                anova_df['Interpretation (alpha=0.05)'] = anova_df['PR(>F)'].apply(interpret_pval)
            add_table_to_slide(slide1, anova_df, left_panel_left, Inches(2.0), left_panel_width, Inches(2.0))
            # Add CV% and CD Value below the ANOVA table
            cv_value = data.get('overall_cv')
            cd_value = data.get('cd_value')
            cv_text = f"Coeff. Variation (CV): {cv_value:.2f}%" if isinstance(cv_value, (int, float)) else "CV: N/A"
            cd_text = f"Critical Diff. (CD): {cd_value:.4f}" if isinstance(cd_value, (int, float)) else "CD: N/A"
            cvcd_shape = slide1.shapes.add_textbox(left_panel_left, Inches(4.1), left_panel_width, Inches(0.5))
            text_frame = cvcd_shape.text_frame
            text_frame.clear()  # Ensure no residual text
            p1 = text_frame.paragraphs[0]
            p1.text = cv_text
            p1.font.size = Pt(10)
            p2 = text_frame.add_paragraph()
            p2.text = cd_text
            p2.font.size = Pt(10)
        # Right panel (Mean separation table)
        right_panel_left = Inches(7.5)
        right_panel_top = Inches(0.5)
        right_panel_width = Inches(5.5)
        mean_separation_results = data.get('mean_separation_results')
        mean_sep_key = None
        if mean_separation_results and isinstance(mean_separation_results, dict) and mean_separation_results.keys():
            mean_sep_key = list(mean_separation_results.keys())[0]
            mean_sep_df = pd.read_json(mean_separation_results[mean_sep_key], orient='records')
            # Sort by first column (A-Z)
            if not mean_sep_df.empty:
                first_col = mean_sep_df.columns[0]
                mean_sep_df = mean_sep_df.sort_values(by=first_col, ascending=True, kind='mergesort')
            mean_sep_title_shape = slide1.shapes.add_textbox(right_panel_left, right_panel_top, right_panel_width, Inches(0.2))
            mean_sep_title_shape.text_frame.text = "Mean Separation Table"
            p = mean_sep_title_shape.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.LEFT
            # Add more space between title and table (move table further down)
            table_top = Inches(0.9)  # was 0.7, now more space
            rows, cols = mean_sep_df.shape[0] + 1, mean_sep_df.shape[1]
            shape = slide1.shapes.add_table(rows, cols, right_panel_left, table_top, right_panel_width, Inches(6.0))
            table = shape.table
            # Set column headers
            for col_index, col_name in enumerate(mean_sep_df.columns):
                cell = table.cell(0, col_index)
                cell.text = str(col_name)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # Add data rows
            for row_index, (idx, row_data) in enumerate(mean_sep_df.iterrows()):
                for col_index, cell_data in enumerate(row_data):
                    cell_text = str(cell_data) if not pd.isna(cell_data) else ""
                    table.cell(row_index + 1, col_index).text = cell_text
            # Adjust font size and row height
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(10)
                # Set row height to max 0.3 inch
                row.height = Inches(0.3)
        # Footer
        footer_text = "Generated using https://vita.chloropy.com; AI can make mistakes, please verify results."
        footer_shape = slide1.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.3))
        footer_shape.text_frame.text = footer_text
        footer_shape.text_frame.paragraphs[0].font.size = Pt(10)
        footer_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # --- Slide 2: Up to 6 charts in 2x3 grid ---
        slide2 = prs.slides.add_slide(blank_slide_layout)
        chart_width = Inches(4.0)  # 3 charts per row
        chart_height = Inches(3.0) # 2 rows
        chart_positions = [
            (Inches(0.3), Inches(0.5)),   # 1st row, 1st col
            (Inches(4.5), Inches(0.5)),   # 1st row, 2nd col
            (Inches(8.7), Inches(0.5)),   # 1st row, 3rd col
            (Inches(0.3), Inches(3.7)),   # 2nd row, 1st col
            (Inches(4.5), Inches(3.7)),   # 2nd row, 2nd col
            (Inches(8.7), Inches(3.7)),   # 2nd row, 3rd col
        ]
        plots = data.get('plots', {})
        chart_keys = [k for k in plots.keys() if plots[k]]
        for idx, key in enumerate(chart_keys[:6]):
            img_b64 = plots[key]
            img_data = base64.b64decode(img_b64)
            img_stream = io.BytesIO(img_data)
            left, top = chart_positions[idx]
            slide2.shapes.add_picture(img_stream, left, top, width=chart_width, height=chart_height)
        # Save presentation to a memory buffer
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        filename = f'{analysis_title.replace(" ", "_")}.pptx'
        return send_file(
            ppt_buffer,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception:
        logging.error("Error generating slide.")
        return jsonify({"error": "An error occurred while generating the slide."}), 500

if __name__ == "__main__":
    debug_mode = os.environ.get("FLASK_ENV", "production") == "development"
    port = int(os.environ.get("PORT", 9006))
    app.run(debug=debug_mode, host="0.0.0.0", port=port)